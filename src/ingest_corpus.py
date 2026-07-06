"""
① 문서 코퍼스 구축 + 오버랩 청킹

코퍼스 구성 (경험8 서술과 일치):
  - corpus/ml_*.txt, ml_*.pdf, ml_*.docx      → 자체 ML 기술보고서
  - corpus/dl_*.txt, dl_*.pptx                → DL 발표자료
  - corpus/public_*.txt, public_*.pdf         → 공개 반도체/설비 문서

지원 포맷: .txt, .pdf, .pptx, .docx
청킹 전략: RecursiveCharacterTextSplitter + chunk_overlap
  → 문서 내 인과관계(문제 인식 → 기법 선택 서술 구조)가 청크 경계에서 끊기지 않도록
    overlap을 충분히 확보 (chunk_size 대비 약 20%).
  → 각 청크에는 원본 파일명, 문서 카테고리(ml/dl/public) metadata를 부여하여
    추후 RAG 답변에서 출처를 구분할 수 있도록 함.
"""

import os
from pathlib import Path
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

CORPUS_DIR = Path(__file__).parent.parent / "corpus"

CHUNK_SIZE = 500
CHUNK_OVERLAP = 100  # chunk_size의 20% -> 문장/문맥 단절 방지


def _category_from_filename(filename: str) -> str:
    """파일명 접두사로 문서 카테고리 태깅 (ml_ / dl_ / public_)"""
    name = filename.lower()
    if name.startswith("ml_"):
        return "ml_technical_report"
    if name.startswith("dl_"):
        return "dl_presentation"
    if name.startswith("public_"):
        return "public_reference"
    return "unclassified"


def load_raw_documents(corpus_dir: Path = CORPUS_DIR) -> list[Document]:
    """corpus/ 폴더의 모든 지원 포맷 문서를 LangChain Document로 로드"""
    docs: list[Document] = []

    for path in sorted(corpus_dir.glob("*")):
        suffix = path.suffix.lower()
        category = _category_from_filename(path.name)

        try:
            if suffix == ".txt":
                text = path.read_text(encoding="utf-8")
                docs.append(
                    Document(
                        page_content=text,
                        metadata={"source": path.name, "category": category},
                    )
                )

            elif suffix == ".pdf":
                from pypdf import PdfReader

                reader = PdfReader(str(path))
                for i, page in enumerate(reader.pages):
                    text = page.extract_text() or ""
                    if text.strip():
                        docs.append(
                            Document(
                                page_content=text,
                                metadata={
                                    "source": path.name,
                                    "category": category,
                                    "page": i + 1,
                                },
                            )
                        )

            elif suffix == ".pptx":
                from pptx import Presentation

                prs = Presentation(str(path))
                for i, slide in enumerate(prs.slides):
                    texts = [
                        shape.text
                        for shape in slide.shapes
                        if shape.has_text_frame and shape.text.strip()
                    ]
                    if texts:
                        docs.append(
                            Document(
                                page_content="\n".join(texts),
                                metadata={
                                    "source": path.name,
                                    "category": category,
                                    "slide": i + 1,
                                },
                            )
                        )

            elif suffix == ".docx":
                import docx

                d = docx.Document(str(path))
                text = "\n".join(p.text for p in d.paragraphs if p.text.strip())
                if text.strip():
                    docs.append(
                        Document(
                            page_content=text,
                            metadata={"source": path.name, "category": category},
                        )
                    )
        except Exception as e:
            print(f"[WARN] {path.name} 로드 실패: {e}")

    return docs


def chunk_documents(
    docs: list[Document],
    chunk_size: int = CHUNK_SIZE,
    chunk_overlap: int = CHUNK_OVERLAP,
) -> list[Document]:
    """오버랩 청킹 적용. 각 청크는 원본 metadata를 상속받는다."""
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    chunks = splitter.split_documents(docs)

    # 각 청크에 chunk_id 부여 (BM25/Chroma 양쪽에서 동일 ID로 참조하기 위함 -> 하이브리드 검색 결과 dedup에 활용)
    for i, chunk in enumerate(chunks):
        chunk.metadata["chunk_id"] = f"{chunk.metadata.get('source', 'unknown')}_{i}"

    return chunks


def build_corpus(corpus_dir: Path = CORPUS_DIR) -> list[Document]:
    raw_docs = load_raw_documents(corpus_dir)
    chunks = chunk_documents(raw_docs)
    return chunks


if __name__ == "__main__":
    chunks = build_corpus()
    print(f"원본 문서 수: {len(load_raw_documents())}")
    print(f"청크 수: {len(chunks)}")
    print("\n--- 샘플 청크 3개 ---")
    for c in chunks[:3]:
        print(f"\n[source={c.metadata['source']}, category={c.metadata['category']}]")
        print(c.page_content[:150].replace("\n", " ") + "...")
